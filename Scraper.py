import requests
import json
from bs4 import BeautifulSoup
import time
import random
import logging
from datetime import datetime
from playwright.async_api import async_playwright
import asyncio
from concurrent.futures import ThreadPoolExecutor
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE  # Added import for MSO_SHAPE
import os

# ==============================================================================
#  CONFIGURATION
# ==============================================================================

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(message)s')
logger = logging.getLogger(__name__)

CONFIG = {
    "MAX_CONCURRENT_REQUESTS": 10,
    "FLARESOLVER_URL": "http://localhost:8191/v1",
    "GEMINI_API_KEY": "Put_your_key_here",  # Replace with your actual Gemini API key
    "GEMINI_MODEL_PREFERENCES": ['models/gemini-2.0-flash', 'models/gemini-2.5-flash', 'models/gemini-1.5-pro-latest'],
    "OUTPUT_FILE": "C:/Users/Hamz/Documents/FT_Analysis_Presentation.pptx",
    "LOGO_PATH": "logo.png",  
    "MAX_RETRIES": 3,
    "REQUEST_TIMEOUT": 90,
    "BATCH_SIZE": 15,  # Process 15 articles per batch to stay within token limits
}

# Configure Gemini API
genai.configure(api_key=CONFIG["GEMINI_API_KEY"])

# Select a supported Gemini model
GEMINI_MODEL = next((m for m in CONFIG["GEMINI_MODEL_PREFERENCES"] if m in [model.name for model in genai.list_models()]), CONFIG["GEMINI_MODEL_PREFERENCES"][0])

# Cache for fetched archive snapshots
fetched_snapshots = {}

# ==============================================================================
#  FLARESOLVER ENGINE
# ==============================================================================

def reset_cookies_and_headers():
    """Set basic headers for FlareSolverr compatibility."""
    return {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
        "Referer": "https://www.google.com/"
    }

def parse_final_article_page(soup: BeautifulSoup, url: str) -> dict:
    """Parses the actual article content from a fully rendered page."""
    title = soup.find('h1').text.strip() if soup.find('h1') else soup.find('title').text.strip() if soup.find('title') else 'Unknown Title'
    body_div = soup.find('div', class_='article__content-body') or soup.find('div', class_='article__content') or \
               soup.find('div', id='CONTENT') or soup.find('article')
    body = body_div.get_text(separator='\n', strip=True) if body_div else 'Content not found'
    return {'title': title, 'link': url, 'content': body}

def fetch_with_flaresolverr(url: str, session_id: str) -> dict:
    """Fetch content using FlareSolverr with retry logic."""
    global fetched_snapshots
    headers = reset_cookies_and_headers()
    for attempt in range(CONFIG["MAX_RETRIES"]):
        try:
            payload = {"cmd": "request.get", "url": url, "maxTimeout": CONFIG["REQUEST_TIMEOUT"] * 1000, "headers": headers, "session": session_id}
            response = requests.post(CONFIG["FLARESOLVER_URL"], json=payload, timeout=CONFIG["REQUEST_TIMEOUT"])
            response.raise_for_status()
            result = response.json()

            if not result.get("solution"):
                raise Exception("FlareSolverr returned no solution.")
            
            html_content = result["solution"]["response"]
            soup = BeautifulSoup(html_content, 'html.parser')

            if "cliché le plus récent" in html_content or "List of URLs" in html_content:
                logger.info(f"Detected archive.ph intermediate page for {url}. Finding latest snapshot link...")
                snapshot_links = soup.find_all('a', href=True)
                latest_snapshot = None
                latest_date = None

                for link in snapshot_links:
                    parent_text = link.find_parent().get_text(strip=True)
                    try:
                        date_obj = datetime.strptime(parent_text, '%d %b %Y %H:%M')
                        if not latest_date or date_obj > latest_date:
                            latest_date = date_obj
                            latest_snapshot = link['href']
                    except ValueError:
                        continue

                if latest_snapshot:
                    if not latest_snapshot.startswith('http'):
                        latest_snapshot = "https://archive.ph" + latest_snapshot
                    if latest_snapshot in fetched_snapshots:
                        logger.info(f"Using cached content for {latest_snapshot}")
                        return fetched_snapshots[latest_snapshot]
                    logger.info(f"Found latest snapshot link: {latest_snapshot}. Fetching final page...")
                    result = fetch_with_flaresolverr(latest_snapshot, session_id)
                    fetched_snapshots[latest_snapshot] = result
                    return result
                else:
                    fallback_link = soup.select_one('div.TEXT-BLOCK a[href^="/"]')
                    if fallback_link:
                        snapshot_url = "https://archive.ph" + fallback_link['href']
                        if snapshot_url in fetched_snapshots:
                            logger.info(f"Using cached content for {snapshot_url}")
                            return fetched_snapshots[snapshot_url]
                        logger.info(f"Found snapshot link via fallback: {snapshot_url}. Fetching final page...")
                        result = fetch_with_flaresolverr(snapshot_url, session_id)
                        fetched_snapshots[snapshot_url] = result
                        return result
                    raise Exception("Could not find a valid snapshot link on the intermediate page.")

            logger.info(f"Retrieved content page for {url}")
            result = parse_final_article_page(soup, url)
            return result

        except requests.exceptions.RequestException as e:
            logger.error(f"FlareSolverr fetch failed for {url} (Attempt {attempt + 1}/{CONFIG['MAX_RETRIES']}): {e}")
            if attempt == CONFIG["MAX_RETRIES"] - 1:
                return {'title': 'Fetch Error', 'link': url, 'content': f'Fetch failed after {CONFIG["MAX_RETRIES"]} attempts: {str(e)}'}
            time.sleep(5)
        except Exception as e:
            logger.error(f"Error processing {url} (Attempt {attempt + 1}/{CONFIG['MAX_RETRIES']}): {e}")
            if attempt == CONFIG["MAX_RETRIES"] - 1:
                return {'title': 'Processing Error', 'link': url, 'content': f'Error after {CONFIG["MAX_RETRIES"]} attempts: {str(e)}'}
            time.sleep(5)

# ==============================================================================
#  PLAYWRIGHT LINK SCRAPER
# ==============================================================================

async def get_all_main_page_links() -> list:
    """Gather all unique article links from FT.com."""
    logger.info("--- Phase 1: Gathering article links from FT.com ---")
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36"
        )
        page = await context.new_page()
        try:
            logger.info("   -> Navigating to https://www.ft.com/...")
            await page.goto("https://www.ft.com/", wait_until='domcontentloaded', timeout=30000)
            logger.info("   -> Handling cookie pop-up...")
            await page.wait_for_timeout(1000)
            for frame in page.frames:
                try:
                    accept_button = frame.get_by_role('button', name='Accept Cookies')
                    if await accept_button.is_visible(timeout=1000):
                        await accept_button.click()
                        await page.wait_for_timeout(1000)
                        break
                except Exception:
                    continue
            
            logger.info("   -> Scraping all unique article links...")
            selector = 'div.js-teaser-headline > a'
            await page.locator(selector).first.wait_for(timeout=10000)
            
            all_locators = await page.locator(selector).all()
            unique_links = [f"https://www.ft.com{await el.get_attribute('href')}" for el in all_locators if await el.get_attribute('href') and (await el.get_attribute('href')).startswith('/content/')]
            
            await browser.close()
            logger.info(f"--- Phase 1 Complete: Found {len(unique_links)} unique article links ---")
            return unique_links
        except Exception as e:
            logger.error(f"Failed to gather links from main page: {e}")
            await browser.close()
            return []

# ==============================================================================
#  PRESENTATION GENERATION
# ==============================================================================
def generate_presentation_batch(articles: list, prs: Presentation = None) -> Presentation:
    """Generate presentation slides for a batch of articles and append to existing presentation if provided."""
    if prs is None:
        logger.info("--- Phase 3: Generating presentation with Gemini ---")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

        prs.slide_master.background.fill.solid()
        prs.slide_master.background.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background

        # Title slide with logo (optional)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Financial Times Article Report"  # Updated to fixed title
        if not title.text_frame.paragraphs:
            title.text_frame.add_paragraph()
        title.text_frame.paragraphs[0].font.size = Pt(28)

        subtitle = slide.placeholders[1]
        subtitle.text = f"Prepared on {datetime.now().strftime('%B %d, %Y')}"
        if not subtitle.text_frame.paragraphs:
            subtitle.text_frame.add_paragraph()
        subtitle.text_frame.paragraphs[0].font.size = Pt(16)

        if CONFIG["LOGO_PATH"] and os.path.exists(CONFIG["LOGO_PATH"]):
            left = Inches(0.5)
            top = Inches(0.5)
            pic = slide.shapes.add_picture(CONFIG["LOGO_PATH"], left, top, height=Inches(1))

    articles_text = json.dumps(articles, ensure_ascii=False)
    
    prompt = f"""
    Given the following JSON data from FT.com articles:
    {articles_text}
    Create a professional presentation outline with 2 slides per article:
    - Slide 1: Article title and a concise introduction (50-70 words) summarizing the article's focus.
    - Slide 2: Conclusion and strategic implications (50-70 words, formal tone).
    Return the output as a JSON object with the structure:
    {{
        "articles": {{
            "article_title": {{
                "slides": [
                    {{"title": "slide_title", "content": "slide_content"}},
                    ...
                ]
            }},
            ...
        }}
    }}
    """

    max_retries = 3
    for attempt in range(max_retries):
        try:
            model = genai.GenerativeModel(GEMINI_MODEL)
            logger.info(f"Sending prompt to Gemini model (Attempt {attempt + 1}/{max_retries})")
            response = model.generate_content(prompt)
            logger.info(f"Gemini response received")
            if response.candidates and response.candidates[0].content and response.candidates[0].content.parts:
                presentation_text = response.candidates[0].content.parts[0].text
                logger.info(f"Raw presentation text: {presentation_text[:500]}...")  # Log first 500 chars for debugging
                if presentation_text.startswith("```json\n") and presentation_text.endswith("\n```"):
                    presentation_text = presentation_text[len("```json\n"):-len("\n```")]
                presentation_data = json.loads(presentation_text)
                break
            else:
                raise ValueError("No valid content in Gemini response")
        except json.JSONDecodeError as e:
            logger.error(f"Error decoding Gemini response (Attempt {attempt + 1}/{max_retries}): {e}. Response text: {presentation_text if 'presentation_text' in locals() else 'No response'}")
            if attempt == max_retries - 1:
                raise
            time.sleep(5)
        except Exception as e:
            logger.error(f"Error generating presentation (Attempt {attempt + 1}/{max_retries}): {e}")
            if attempt == max_retries - 1:
                raise
            time.sleep(5)

    # Article slides with adjusted layout
    layout = prs.slide_layouts[1]  # Title and Content layout
    for article_title, data in presentation_data["articles"].items():
        for slide_data in data["slides"]:
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            logger.info(f"Processing slide: {slide_data['title']}")
            title.text = slide_data["title"]
            if not title.text_frame.paragraphs:
                title.text_frame.add_paragraph()
            title.text_frame.paragraphs[0].font.size = Pt(20)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)  # Dark blue

            body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
            logger.info(f"Processing body for: {slide_data['title']}")
            text_frame = body.text_frame
            text_frame.clear()

            content = slide_data["content"] if isinstance(slide_data["content"], str) else '\n'.join(slide_data["content"])
            p = text_frame.add_paragraph()
            p.text = content

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(64, 64, 64)  # Dark gray

            if slide_data == data["slides"][-1]:  # Add line for the last slide of each article
                left = Inches(0.5)
                top = Inches(4.5)
                try:
                    logger.info(f"Adding rectangle shape to slide: {slide_data['title']}")
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(9), Inches(0.5))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0, 32, 96)
                    shape.line.color.rgb = RGBColor(0, 32, 96)
                except Exception as e:
                    logger.error(f"Failed to add shape to slide {slide_data['title']}: {e}")

    return prs
# ==============================================================================
#  MAIN EXECUTION WORKFLOW
# ==============================================================================

async def process_article(url: str, session_id: str, semaphore) -> dict:
    """Process a single article using FlareSolverr."""
    async with semaphore:
        loop = asyncio.get_running_loop()
        article = await loop.run_in_executor(None, fetch_with_flaresolverr, url, session_id)
        time.sleep(random.uniform(0.5, 1))
        return article

async def main() -> None:
    """Main execution workflow for scraping and presentation generation."""
    article_links = await get_all_main_page_links()
    if not article_links:
        logger.error("Stopping script as no links were found.")
        return

    logger.info(f"\n--- Phase 2: Processing {len(article_links)} links with FlareSolverr ---")
    session_id = f"session_{random.randint(10000, 99999)}"
    semaphore = asyncio.Semaphore(CONFIG["MAX_CONCURRENT_REQUESTS"])
    
    tasks = [process_article(f"https://archive.ph/{link}", session_id, semaphore) for link in article_links]
    scraped_articles = await asyncio.gather(*tasks)
    
    valid_articles = [a for a in scraped_articles if a and a.get('title') != 'Fetch Error' and a.get('title') != 'Cached Content' and a.get('title') != 'Processing Error']
    with open('bypassed_ft_articles.json', 'w', encoding='utf-8') as f:
        json.dump(valid_articles, f, ensure_ascii=False, indent=4)
    logger.info(f"--- Phase 2 Complete: Saved {len(valid_articles)} processed articles ---")

    with open('bypassed_ft_articles.json', 'r', encoding='utf-8') as f:
        articles = json.load(f)
    
    # Initialize presentation
    prs = None
    # Process articles in batches
    for i in range(0, len(articles), CONFIG["BATCH_SIZE"]):
        batch_articles = articles[i:i + CONFIG["BATCH_SIZE"]]
        prs = generate_presentation_batch(batch_articles, prs)

    # Save the single presentation with error handling
    if prs:
        try:
            prs.save(CONFIG["OUTPUT_FILE"])
            logger.info(f"--- Phase 3 Complete: Presentation saved as {CONFIG['OUTPUT_FILE']} ---")
        except PermissionError as pe:
            logger.error(f"Permission denied saving presentation: {pe}. Check file access or close open instances.")
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
    else:
        logger.error("No presentation generated.")

if __name__ == "__main__":
    asyncio.run(main())